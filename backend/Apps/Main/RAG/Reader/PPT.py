from pptx import Presentation

from backend.Apps.Main.RAG.Dataclass import FileInfo
from backend.Lib.Error import FileNotSupported
from .Base import BaseRAG

class PPT(BaseRAG):
  @classmethod
  def is_document(cls, file_info: FileInfo) -> bool:
    is_correct_content_type = file_info.content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    return file_info.filename.lower().endswith(".pptx") or is_correct_content_type
  
  @classmethod
  def read(cls, file_info: FileInfo) -> str:
    try:
      file_info.stream.seek(0)
      prs = Presentation(file_info.stream)
      texts = []
      for slide in prs.slides:
          for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip(): # type: ignore
                texts.append(shape.text) # type: ignore
      return "\n".join(texts) # type: ignore
    except Exception as e:
      raise FileNotSupported(str(e))